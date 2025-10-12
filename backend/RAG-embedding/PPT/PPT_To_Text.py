# ppt_to_text.py

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
import io
import os

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\tesseract.exe"  

def ppt_to_text(ppt_file, output_dir="Text_files", output_name=None):
    """
    Extract all text from a PPT/PPTX file including:
    - Normal text boxes
    - Tables
    - Charts
    - Text inside images (OCR)
    
    Saves output as a .txt file with the same base name as the PPT file.
    """
    presentation = Presentation(ppt_file)
    ppt_text = ""

    for slide_num, slide in enumerate(presentation.slides, start=1):
        ppt_text += f"--- Slide {slide_num} ---\n"

        for shape in slide.shapes:
            # 1. Normal text
            if hasattr(shape, "text") and shape.text.strip():
                ppt_text += "[Text Box]:\n" + shape.text + "\n"

            # 2. Tables
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                ppt_text += "[Table]:\n"
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    ppt_text += "\t".join(row_text) + "\n"

            # 3. Charts (extract text from chart series and categories)
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                ppt_text += "[Chart]:\n"
                chart = shape.chart
                for series in chart.series:
                    ppt_text += f"Series: {series.name}\n"
                    ppt_text += "Values: " + ", ".join([str(v) for v in series.values]) + "\n"
                categories = [str(c) for c in chart.chart_data.categories]
                ppt_text += "Categories: " + ", ".join(categories) + "\n"

            # 4. Images (OCR)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_stream = io.BytesIO(image_bytes)
                img = Image.open(image_stream)
                text_in_image = pytesseract.image_to_string(img)
                if text_in_image.strip():
                    ppt_text += "[Text in Image]:\n" + text_in_image + "\n"

        ppt_text += "\n"

    # Save output with specified name or original base name
    os.makedirs(output_dir, exist_ok=True)
    base_name = output_name if output_name else os.path.splitext(os.path.basename(ppt_file))[0]
    output_file = os.path.join(output_dir, f"{base_name}.txt")
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(ppt_text)

    print(f"Extraction complete! Saved to {output_file}")
    return output_file
